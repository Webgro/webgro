"""
Build the Ithaca Roofing proposal deck.

Design system — "cream studio" variant:
- Cover is dark ink for drama; every other slide is wg-mist cream (#F4F6FB)
  with pure-white cards. Reads closer to Koto / Pentagram than agency
  dark-mode theatre.
- wg-blue #2D8DFF, wg-violet #7C3AED, wg-teal #00C9A7 as the triad accent,
  deployed as bars, dots, and medallions rather than as whole-slide washes.
- Space Grotesk for display type (geometric, flared terminals, plenty of
  character) with Inter as the body sibling. JetBrains Mono carries the
  uppercase tracked eyebrow labels. All three are bundled in Google Slides.

Google Slides imports .pptx natively: upload to Drive, right-click,
Open with Google Slides. Fonts map cleanly.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

# -----------------------------------------------------------------------------
# Palette — light (cream) variant. The wg-mist token from the site becomes the
# page background; text flips to near-black ink; cards are pure white so they
# catch the eye as "raised". The blue/violet/teal accent triad is unchanged,
# so the proposal still reads "Webgro" without being dark-mode theatre.
# -----------------------------------------------------------------------------
INK = RGBColor(0x07, 0x08, 0x0C)          # primary text (+ cover bg)
INK_RAISED = RGBColor(0xFF, 0xFF, 0xFF)   # raised cards (white on cream)
FOREGROUND = INK                          # primary text alias (content slides)
PAGE_BG = RGBColor(0xF4, 0xF6, 0xFB)      # wg-mist cream page background
COVER_BG = INK                            # cover slide stays dark for drama
COVER_FG = RGBColor(0xED, 0xEF, 0xF4)     # ink-slide headings
COVER_SLATE = RGBColor(0x9A, 0xA6, 0xBC)  # ink-slide body/subtitles
BLUE = RGBColor(0x2D, 0x8D, 0xFF)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
TEAL = RGBColor(0x00, 0xC9, 0xA7)
SLATE = RGBColor(0x8B, 0x94, 0xA8)        # tertiary/mono labels (lighter for cream bg)
SLATE_SOFT = RGBColor(0x52, 0x5B, 0x6E)   # secondary body copy (darker for legibility)
WHITE_20 = RGBColor(0xE1, 0xE5, 0xEF)     # border mid
WHITE_10 = RGBColor(0xE6, 0xEA, 0xF2)     # hairline border on cards

DISPLAY_FONT = "Space Grotesk"
MONO_FONT = "JetBrains Mono"
BODY_FONT = "Inter"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

IMG_DIR = Path("/Users/michaelbroadbridge/Desktop/Claude/webgro-site/public/work")
BRAND_DIR = Path("/Users/michaelbroadbridge/Desktop/Claude/webgro-site/public/brand")
OUT_PATH = Path(
    "/Users/michaelbroadbridge/Desktop/Claude/webgro-site/proposals/ithaca-roofing-proposal.pptx"
)

# -----------------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------------


def set_slide_bg(slide, rgb: RGBColor) -> None:
    """Paint the whole slide a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
    """Draw a filled rectangle with an optional border."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_rounded_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, radius=0.08):
    """Rounded rectangle with configurable corner radius (0-0.5 of short side)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_gradient_rect(slide, x, y, w, h, stops):
    """Horizontal gradient rectangle. stops = [(pos, RGBColor), …] with pos 0..100000."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    shp.shadow.inherit = False
    # Clear solid fill, drop in a gradientFill in xml
    spPr = shp.fill._xPr.find(qn("a:spPr")) if False else None  # unused path
    sp = shp._element
    # Replace fill node
    fill_parent = sp.find(".//" + qn("p:spPr"))
    if fill_parent is None:
        fill_parent = sp.find(qn("p:spPr"))
    # Find and drop existing solidFill or noFill
    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
        existing = fill_parent.find(qn(tag))
        if existing is not None:
            fill_parent.remove(existing)
    gradFill = etree.SubElement(fill_parent, qn("a:gradFill"))
    gradFill.set("flip", "none")
    gradFill.set("rotWithShape", "1")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, color in stops:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", f"{color.rgb if hasattr(color, 'rgb') else color}")
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", "0")  # horizontal
    lin.set("scaled", "1")
    tileRect = etree.SubElement(gradFill, qn("a:tileRect"))
    # Ensure gradFill is placed BEFORE the line (a:ln) element
    ln = fill_parent.find(qn("a:ln"))
    if ln is not None:
        # move gradFill just before ln
        fill_parent.remove(gradFill)
        ln.addprevious(gradFill)
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font=BODY_FONT,
    size=14,
    bold=False,
    italic=False,
    color=FOREGROUND,
    align=PP_ALIGN.LEFT,
    line_spacing=1.2,
    anchor=MSO_ANCHOR.TOP,
    tracking=None,
):
    """Add a text box. Returns the textframe for further tweaking."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    # First paragraph already exists
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if tracking is not None:
            # Track-out in 1/100 of a point; python-pptx doesn't expose spc
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(tracking)))
    return tf


def add_eyebrow(slide, x, y, w, text, color=BLUE):
    """Small uppercase mono label with letter-spacing, used throughout the deck."""
    return add_text(
        slide,
        x,
        y,
        w,
        Inches(0.3),
        text.upper(),
        font=MONO_FONT,
        size=10,
        color=color,
        tracking=280,  # ~0.25em
    )


def add_dot_accent(slide, x, y, color):
    """Small pill accent like the website's section eyebrow dashes."""
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.4), Inches(0.04)
    )
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    pill.shadow.inherit = False
    return pill


def add_image(slide, path: Path, x, y, w, h, *, cover=True):
    """Add an image. If cover=True we fit-to-fill and crop to the frame."""
    if not path.exists():
        # Fallback placeholder
        add_rect(slide, x, y, w, h, INK_RAISED, WHITE_10)
        return None
    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if cover:
        # Compute crop so image fills frame, preserving aspect ratio.
        # PIL needed for natural size.
        from PIL import Image as PImage
        with PImage.open(path) as im:
            iw, ih = im.size
        frame_ratio = w / h
        img_ratio = iw / ih
        # If image is wider than frame, crop left/right; if taller, crop top/bottom
        if img_ratio > frame_ratio:
            # crop horizontally
            new_w = ih * frame_ratio
            crop_x = (iw - new_w) / 2
            pic.crop_left = crop_x / iw
            pic.crop_right = crop_x / iw
        else:
            # crop vertically
            new_h = iw / frame_ratio
            crop_y = (ih - new_h) / 2
            pic.crop_top = crop_y / ih
            pic.crop_bottom = crop_y / ih
    return pic


def add_page_number(slide, n: int, total: int):
    """Top-right mono page indicator like the site's [ 01 ] eyebrow."""
    add_text(
        slide,
        SLIDE_W - Inches(1.4),
        SLIDE_H - Inches(0.45),
        Inches(1.2),
        Inches(0.3),
        f"[ {n:02d} / {total:02d} ]",
        font=MONO_FONT,
        size=9,
        color=SLATE,
        tracking=200,
        align=PP_ALIGN.RIGHT,
    )


def add_footer(slide):
    """Webgro mark bottom-left on every slide (except cover)."""
    add_text(
        slide,
        Inches(0.6),
        SLIDE_H - Inches(0.45),
        Inches(3.0),
        Inches(0.3),
        "webgro.co.uk",
        font=MONO_FONT,
        size=9,
        color=SLATE,
        tracking=200,
    )


# -----------------------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------------------


def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, COVER_BG)

    # Ambient gradient glow, top-right quadrant
    for offset, alpha in [(0.0, 0.35), (0.4, 0.18), (1.0, 0.0)]:
        pass  # python-pptx's radial gradients are awkward; skip, use accent bar instead

    # Top eyebrow
    add_eyebrow(slide, Inches(0.6), Inches(0.6), Inches(4), "[ Webgro / 2026 ] Proposal")

    # Gradient accent bar (blue → violet → teal)
    add_gradient_rect(
        slide,
        Inches(0.6),
        Inches(1.15),
        Inches(4.5),
        Inches(0.08),
        [(0, BLUE), (50000, VIOLET), (100000, TEAL)],
    )

    # Big title
    add_text(
        slide,
        Inches(0.6),
        Inches(1.4),
        Inches(11),
        Inches(3.5),
        "A proposal for\nIthaca Roofing.",
        font=DISPLAY_FONT,
        size=80,
        bold=True,
        line_spacing=0.95,
        color=COVER_FG,
    )

    # Lead / subhead
    add_text(
        slide,
        Inches(0.6),
        Inches(4.8),
        Inches(9.5),
        Inches(1.3),
        "A custom WordPress rebuild of the existing site, with drag-and-drop "
        "editing and ongoing maintenance. Written April 2026.",
        font=BODY_FONT,
        size=20,
        color=COVER_SLATE,
        line_spacing=1.35,
    )

    # Webgro logo bottom-left (replaces the wordmark text).
    # Source PNG is 450 x 146, ~3.08 aspect; render at Inches(2.3) wide so
    # it sits at roughly the same visual weight as the old text wordmark.
    logo_w = Inches(2.3)
    logo_h = Inches(2.3 / (450 / 146))
    slide.shapes.add_picture(
        str(BRAND_DIR / "logo-white.png"),
        Inches(0.6),
        SLIDE_H - Inches(1.3),
        width=logo_w,
        height=logo_h,
    )
    add_text(
        slide,
        Inches(0.6),
        SLIDE_H - Inches(0.55),
        Inches(10),
        Inches(0.3),
        "AI-FIRST eCOMMERCE & WORDPRESS AGENCY  ·  BRACKNELL, UK",
        font=MONO_FONT,
        size=9,
        color=COVER_SLATE,
        tracking=260,
    )

    # Accent dot bottom-right — a true dot (circle)
    dot_d = Inches(0.16)
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        SLIDE_W - Inches(0.6) - dot_d,
        SLIDE_H - Inches(0.85),
        dot_d,
        dot_d,
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = TEAL
    dot.line.fill.background()
    dot.shadow.inherit = False


def build_section_divider(prs, num: str, title: str, sub: str, accent: RGBColor, n: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAGE_BG)

    # Left vertical accent bar
    add_rect(slide, Inches(0.6), Inches(2.6), Inches(0.08), Inches(2.3), accent)

    # Big section number
    add_text(
        slide,
        Inches(0.95),
        Inches(2.5),
        Inches(3),
        Inches(1.2),
        num,
        font=DISPLAY_FONT,
        size=120,
        bold=True,
        color=accent,
        line_spacing=0.9,
    )

    # Title
    add_text(
        slide,
        Inches(0.95),
        Inches(3.9),
        Inches(11.5),
        Inches(1.5),
        title,
        font=DISPLAY_FONT,
        size=64,
        bold=True,
        color=FOREGROUND,
        line_spacing=1.0,
    )

    # Subtitle
    add_text(
        slide,
        Inches(0.95),
        Inches(5.15),
        Inches(10),
        Inches(0.8),
        sub,
        font=BODY_FONT,
        size=20,
        color=SLATE_SOFT,
        line_spacing=1.35,
    )

    add_footer(slide)
    add_page_number(slide, n, total)


def build_hello(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAGE_BG)

    # Left column
    add_eyebrow(slide, Inches(0.6), Inches(0.75), Inches(4), "[ 01 ] Who we are", color=BLUE)

    add_text(
        slide,
        Inches(0.6),
        Inches(1.3),
        Inches(7),
        Inches(2.5),
        "Hello.\nWe're Webgro.",
        font=DISPLAY_FONT,
        size=72,
        bold=True,
        color=FOREGROUND,
        line_spacing=0.95,
    )

    add_text(
        slide,
        Inches(0.6),
        Inches(4.4),
        Inches(7),
        Inches(2.3),
        "A senior studio building the sharp end of modern web. Fifteen years deep "
        "on WordPress and Shopify, with custom design, development, and ongoing "
        "care built in. No template bloat, no agency-hopping between teams, "
        "no hero launches that crumble the week after.",
        font=BODY_FONT,
        size=16,
        color=SLATE_SOFT,
        line_spacing=1.5,
    )

    # Right column — three stat callouts in a stacked card
    card_x = Inches(8.2)
    card_y = Inches(1.2)
    card_w = Inches(4.6)
    card_h = Inches(5.4)
    add_rounded_rect(slide, card_x, card_y, card_w, card_h, INK_RAISED, WHITE_10, radius=0.04)

    stats = [
        ("15+ YRS", "WordPress & eCommerce", BLUE),
        ("5x", "Award-winning studio", VIOLET),
        ("120+", "Sites launched & maintained", TEAL),
    ]
    for i, (value, label, color) in enumerate(stats):
        yy = card_y + Inches(0.5 + i * 1.7)
        add_text(
            slide,
            card_x + Inches(0.4),
            yy,
            card_w - Inches(0.8),
            Inches(0.9),
            value,
            font=DISPLAY_FONT,
            size=54,
            bold=True,
            color=color,
        )
        add_text(
            slide,
            card_x + Inches(0.4),
            yy + Inches(0.95),
            card_w - Inches(0.8),
            Inches(0.5),
            label,
            font=BODY_FONT,
            size=13,
            color=SLATE_SOFT,
        )
        # separator (except last)
        if i < len(stats) - 1:
            add_rect(
                slide,
                card_x + Inches(0.4),
                yy + Inches(1.45),
                card_w - Inches(0.8),
                Emu(9525),  # 1pt
                WHITE_10,
            )

    add_footer(slide)
    add_page_number(slide, n, total)


def build_capabilities(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAGE_BG)

    add_eyebrow(slide, Inches(0.6), Inches(0.6), Inches(5), "[ 02 ] What we do")
    add_text(
        slide,
        Inches(0.6),
        Inches(1.05),
        Inches(12),
        Inches(1.3),
        "Six capabilities, one studio.",
        font=DISPLAY_FONT,
        size=54,
        bold=True,
        color=FOREGROUND,
        line_spacing=1.0,
    )
    add_text(
        slide,
        Inches(0.6),
        Inches(2.15),
        Inches(10),
        Inches(0.6),
        "No agency-hopping for strategy, build, or marketing. Same team across the whole loop.",
        font=BODY_FONT,
        size=15,
        color=SLATE_SOFT,
    )

    # 3 x 2 grid of capability cards
    caps = [
        ("01", "Websites", "Shopify & WordPress builds, fast by default.", BLUE),
        ("02", "Consultancy", "Senior eCommerce expertise, shaped to your team.", VIOLET),
        ("03", "Automation & AI", "Production AI that earns its keep week one.", TEAL),
        ("04", "SEO", "Technical foundations, content that compounds.", BLUE),
        ("05", "Marketing", "Paid + email + analytics as one loop.", VIOLET),
        ("06", "Design", "Systems, not decoration.", TEAL),
    ]

    grid_x = Inches(0.6)
    grid_y = Inches(3.0)
    col_w = Inches(4.04)
    row_h = Inches(2.05)
    gap = Inches(0.12)

    for i, (num, name, body, color) in enumerate(caps):
        col = i % 3
        row = i // 3
        x = grid_x + col * (col_w + gap)
        y = grid_y + row * (row_h + gap)
        add_rounded_rect(slide, x, y, col_w, row_h, INK_RAISED, WHITE_10, radius=0.06)
        # colored accent disc
        disc = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.35), y + Inches(0.35), Inches(0.18), Inches(0.18)
        )
        disc.fill.solid()
        disc.fill.fore_color.rgb = color
        disc.line.fill.background()
        disc.shadow.inherit = False
        # number
        add_text(
            slide,
            x + Inches(0.65),
            y + Inches(0.3),
            Inches(1),
            Inches(0.3),
            num,
            font=MONO_FONT,
            size=10,
            color=color,
            tracking=200,
        )
        # name
        add_text(
            slide,
            x + Inches(0.35),
            y + Inches(0.75),
            col_w - Inches(0.7),
            Inches(0.6),
            name,
            font=DISPLAY_FONT,
            size=24,
            bold=True,
            color=FOREGROUND,
        )
        # body
        add_text(
            slide,
            x + Inches(0.35),
            y + Inches(1.3),
            col_w - Inches(0.7),
            Inches(0.9),
            body,
            font=BODY_FONT,
            size=12,
            color=SLATE_SOFT,
            line_spacing=1.35,
        )

    add_footer(slide)
    add_page_number(slide, n, total)


def build_awards(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAGE_BG)

    add_eyebrow(slide, Inches(0.6), Inches(0.6), Inches(5), "[ 03 ] Awards")
    add_text(
        slide,
        Inches(0.6),
        Inches(1.05),
        Inches(12),
        Inches(1.3),
        "Five-time award winners.",
        font=DISPLAY_FONT,
        size=54,
        bold=True,
        color=FOREGROUND,
        line_spacing=1.0,
    )

    awards = [
        ("2024", "Best Web Design Agency", "South East England", VIOLET),
        ("2024", "Best eCommerce Consultant", "South East England", TEAL),
        ("2022", "Best Web Design Agency", "United Kingdom", BLUE),
        ("2021", "Best Web Design Agency", "South East England", BLUE),
        ("2020", "Best Web Design Agency", "Berkshire", BLUE),
    ]

    card_w = Inches(2.4)
    card_h = Inches(3.2)
    gap = Inches(0.15)
    total_w = 5 * card_w + 4 * gap
    start_x = (SLIDE_W - total_w) / 2
    card_y = Inches(2.85)

    for i, (year, title_, region, color) in enumerate(awards):
        x = start_x + i * (card_w + gap)
        add_rounded_rect(slide, x, card_y, card_w, card_h, INK_RAISED, WHITE_10, radius=0.04)
        # Medallion icon: colored disc with a white 5-point star inside
        disc_d = Inches(1.1)
        disc_x = x + (card_w - disc_d) / 2
        disc_y = card_y + Inches(0.5)
        disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, disc_x, disc_y, disc_d, disc_d)
        disc.fill.solid()
        disc.fill.fore_color.rgb = color
        disc.line.fill.background()
        disc.shadow.inherit = False
        # Star inset
        star_d = Inches(0.55)
        star_x = disc_x + (disc_d - star_d) / 2
        star_y = disc_y + (disc_d - star_d) / 2
        star = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, star_x, star_y, star_d, star_d)
        star.fill.solid()
        star.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        star.line.fill.background()
        star.shadow.inherit = False
        # year
        add_text(
            slide,
            x + Inches(0.2),
            card_y + Inches(1.7),
            card_w - Inches(0.4),
            Inches(0.55),
            year,
            font=DISPLAY_FONT,
            size=32,
            bold=True,
            color=color,
            align=PP_ALIGN.CENTER,
        )
        # title
        add_text(
            slide,
            x + Inches(0.2),
            card_y + Inches(2.25),
            card_w - Inches(0.4),
            Inches(0.6),
            title_,
            font=BODY_FONT,
            size=12,
            bold=True,
            color=FOREGROUND,
            align=PP_ALIGN.CENTER,
            line_spacing=1.2,
        )
        # region
        add_text(
            slide,
            x + Inches(0.2),
            card_y + Inches(2.75),
            card_w - Inches(0.4),
            Inches(0.4),
            region.upper(),
            font=MONO_FONT,
            size=9,
            color=SLATE,
            align=PP_ALIGN.CENTER,
            tracking=200,
        )

    add_footer(slide)
    add_page_number(slide, n, total)


def build_case_study(
    prs,
    n,
    total,
    *,
    client: str,
    url: str,
    tagline: str,
    services: str,
    image_path: Path,
    accent: RGBColor,
    year: str,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAGE_BG)

    # Left text column.
    # Client names vary in length ("Fun Cases" → "Space 4 U Self Storage").
    # We drop the display size when the name is long enough to wrap, so the
    # title stays single-line and the URL badge below doesn't collide.
    title_size = 60 if len(client) <= 14 else 44
    # Reserve room for 2 lines anyway, so very long names (if they still wrap
    # at 44pt) don't collide with what follows.
    title_h = Inches(1.8)

    add_eyebrow(
        slide,
        Inches(0.6),
        Inches(0.7),
        Inches(6),
        f"[ Case study / {year} ]",
        color=accent,
    )

    add_text(
        slide,
        Inches(0.6),
        Inches(1.2),
        Inches(7),
        title_h,
        client,
        font=DISPLAY_FONT,
        size=title_size,
        bold=True,
        color=FOREGROUND,
        line_spacing=1.0,
    )

    # url badge — pushed below the reserved title block
    url_y = Inches(3.05)
    url_h = Inches(0.42)
    # Measure-less badge: fixed-width rounded pill
    add_rounded_rect(
        slide, Inches(0.6), url_y, Inches(3.8), url_h, INK_RAISED, WHITE_10, radius=0.5
    )
    # url dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.85),
        url_y + Inches(0.16),
        Inches(0.1),
        Inches(0.1),
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent
    dot.line.fill.background()
    dot.shadow.inherit = False
    add_text(
        slide,
        Inches(1.05),
        url_y + Inches(0.08),
        Inches(3.4),
        url_h - Inches(0.15),
        url,
        font=MONO_FONT,
        size=11,
        color=FOREGROUND,
        tracking=100,
    )

    # tagline (pushed down to clear the taller title + URL badge)
    add_text(
        slide,
        Inches(0.6),
        Inches(3.85),
        Inches(6.8),
        Inches(2.1),
        tagline,
        font=DISPLAY_FONT,
        size=22,
        color=SLATE_SOFT,
        line_spacing=1.35,
    )

    # services line
    add_text(
        slide,
        Inches(0.6),
        Inches(6.15),
        Inches(7),
        Inches(0.4),
        "SERVICES",
        font=MONO_FONT,
        size=9,
        color=SLATE,
        tracking=220,
    )
    add_text(
        slide,
        Inches(0.6),
        Inches(6.45),
        Inches(7),
        Inches(0.5),
        services,
        font=BODY_FONT,
        size=14,
        color=FOREGROUND,
    )

    # Right image — big rounded rectangle "browser card"
    img_x = Inches(7.7)
    img_y = Inches(1.0)
    img_w = Inches(5.1)
    img_h = Inches(5.5)
    # card background (acts as border/backdrop)
    add_rounded_rect(slide, img_x, img_y, img_w, img_h, INK_RAISED, WHITE_10, radius=0.03)
    # image inset
    inset = Inches(0.12)
    add_image(
        slide,
        image_path,
        img_x + inset,
        img_y + inset,
        img_w - 2 * inset,
        img_h - 2 * inset,
        cover=True,
    )

    # accent bar bottom of image card
    add_rect(
        slide,
        img_x + Inches(0.12),
        img_y + img_h - Inches(0.12),
        Inches(0.8),
        Emu(25400),  # ~2pt
        accent,
    )

    add_footer(slide)
    add_page_number(slide, n, total)


def build_what_we_provide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAGE_BG)

    add_eyebrow(slide, Inches(0.6), Inches(0.6), Inches(5), "[ Scope ] What you get")
    add_text(
        slide,
        Inches(0.6),
        Inches(1.05),
        Inches(12),
        Inches(1.3),
        "Three lines of delivery.",
        font=DISPLAY_FONT,
        size=54,
        bold=True,
        color=FOREGROUND,
        line_spacing=1.0,
    )

    items = [
        (
            "01",
            "Custom WordPress design & build",
            "A fresh design tailored to Ithaca Roofing, then built from the ground up on WordPress. No third-party themes, no template bloat. Every section shaped around your content and your customers.",
            BLUE,
        ),
        (
            "02",
            "Drag-and-drop editing",
            "Once the build's signed off, we convert every section into editable blocks. Your team can update copy, swap images, and add new pages without touching code or calling a developer.",
            VIOLET,
        ),
        (
            "03",
            "Ongoing maintenance",
            "Managed hosting, daily backups, weekly CMS and plugin updates, plus performance and security checks. Your site stays fast, secure, and current, month after month.",
            TEAL,
        ),
    ]
    row_y = Inches(2.7)
    row_h = Inches(1.35)
    gap = Inches(0.15)
    for i, (num, name, body, color) in enumerate(items):
        y = row_y + i * (row_h + gap)
        add_rounded_rect(slide, Inches(0.6), y, Inches(12.13), row_h, INK_RAISED, WHITE_10, radius=0.06)
        # left colored pip
        add_rect(slide, Inches(0.6), y, Inches(0.08), row_h, color)
        # number
        add_text(
            slide,
            Inches(0.95),
            y + Inches(0.35),
            Inches(0.6),
            Inches(0.4),
            num,
            font=MONO_FONT,
            size=11,
            color=color,
            tracking=200,
        )
        # name
        add_text(
            slide,
            Inches(1.8),
            y + Inches(0.25),
            Inches(4.5),
            Inches(0.6),
            name,
            font=DISPLAY_FONT,
            size=22,
            bold=True,
            color=FOREGROUND,
        )
        # body
        add_text(
            slide,
            Inches(6.5),
            y + Inches(0.3),
            Inches(6.2),
            Inches(1.0),
            body,
            font=BODY_FONT,
            size=13,
            color=SLATE_SOFT,
            line_spacing=1.4,
        )

    add_footer(slide)
    add_page_number(slide, n, total)


def build_investment(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAGE_BG)

    add_eyebrow(slide, Inches(0.6), Inches(0.6), Inches(5), "[ Investment ] For Ithaca Roofing")
    add_text(
        slide,
        Inches(0.6),
        Inches(1.05),
        Inches(12),
        Inches(1),
        "The numbers.",
        font=DISPLAY_FONT,
        size=46,
        bold=True,
        color=FOREGROUND,
    )

    # Big number card
    big_x = Inches(0.6)
    big_y = Inches(2.3)
    big_w = Inches(7.2)
    big_h = Inches(4.2)
    add_rounded_rect(slide, big_x, big_y, big_w, big_h, INK_RAISED, WHITE_10, radius=0.05)

    add_text(
        slide,
        big_x + Inches(0.5),
        big_y + Inches(0.45),
        Inches(4),
        Inches(0.4),
        "ONE-OFF BUILD",
        font=MONO_FONT,
        size=10,
        color=BLUE,
        tracking=260,
    )
    # huge £
    add_text(
        slide,
        big_x + Inches(0.5),
        big_y + Inches(1.0),
        big_w - Inches(1),
        Inches(2.0),
        "£3,750",
        font=DISPLAY_FONT,
        size=130,
        bold=True,
        color=FOREGROUND,
        line_spacing=0.95,
    )
    add_text(
        slide,
        big_x + Inches(0.5),
        big_y + Inches(3.1),
        big_w - Inches(1),
        Inches(1),
        "Custom WordPress redesign and build, converted to drag-and-drop editing. Delivered, deployed, yours.",
        font=BODY_FONT,
        size=15,
        color=SLATE_SOFT,
        line_spacing=1.4,
    )

    # Right side — day rate card
    day_x = Inches(8.0)
    day_y = Inches(2.3)
    day_w = Inches(4.7)
    day_h = Inches(4.2)
    add_rounded_rect(slide, day_x, day_y, day_w, day_h, INK_RAISED, WHITE_10, radius=0.05)
    add_text(
        slide,
        day_x + Inches(0.45),
        day_y + Inches(0.45),
        Inches(4),
        Inches(0.4),
        "ONGOING MAINTENANCE",
        font=MONO_FONT,
        size=10,
        color=VIOLET,
        tracking=260,
    )
    add_text(
        slide,
        day_x + Inches(0.45),
        day_y + Inches(1.0),
        day_w - Inches(0.9),
        Inches(1.4),
        "£120",
        font=DISPLAY_FONT,
        size=76,
        bold=True,
        color=FOREGROUND,
        line_spacing=0.95,
    )
    add_text(
        slide,
        day_x + Inches(0.45),
        day_y + Inches(2.15),
        day_w - Inches(0.9),
        Inches(0.4),
        "per month",
        font=BODY_FONT,
        size=16,
        color=SLATE_SOFT,
    )

    # body
    add_text(
        slide,
        day_x + Inches(0.45),
        day_y + Inches(2.85),
        day_w - Inches(0.9),
        Inches(1.3),
        "Managed hosting, daily backups, weekly CMS and plugin updates, "
        "plus performance and security checks. No lock-in, cancel anytime.",
        font=BODY_FONT,
        size=12,
        color=SLATE_SOFT,
        line_spacing=1.4,
    )

    add_footer(slide)
    add_page_number(slide, n, total)


def build_ready(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAGE_BG)

    add_eyebrow(slide, Inches(0.6), Inches(0.6), Inches(5), "[ Next ] Ready to Gro?")

    # Big gradient-filled heading (we fake a gradient by drawing a white title
    # and a second accent-coloured phrase underneath).
    add_text(
        slide,
        Inches(0.6),
        Inches(1.15),
        Inches(11),
        Inches(1.7),
        "Let's build it.",
        font=DISPLAY_FONT,
        size=88,
        bold=True,
        color=FOREGROUND,
        line_spacing=1.0,
    )
    # italic accent line
    add_text(
        slide,
        Inches(0.6),
        Inches(2.45),
        Inches(11),
        Inches(1.4),
        "Say go and we're rolling.",
        font=DISPLAY_FONT,
        size=52,
        italic=True,
        color=TEAL,
        line_spacing=1.0,
    )

    # Contact card (left)
    card_x = Inches(0.6)
    card_y = Inches(4.3)
    card_w = Inches(6)
    card_h = Inches(2.6)
    add_rounded_rect(slide, card_x, card_y, card_w, card_h, INK_RAISED, WHITE_10, radius=0.05)
    add_text(
        slide,
        card_x + Inches(0.45),
        card_y + Inches(0.35),
        Inches(3),
        Inches(0.35),
        "TALK TO US",
        font=MONO_FONT,
        size=10,
        color=BLUE,
        tracking=260,
    )
    add_text(
        slide,
        card_x + Inches(0.45),
        card_y + Inches(0.8),
        card_w - Inches(0.9),
        Inches(0.5),
        "hello@webgro.co.uk",
        font=DISPLAY_FONT,
        size=22,
        bold=True,
        color=FOREGROUND,
    )
    add_text(
        slide,
        card_x + Inches(0.45),
        card_y + Inches(1.4),
        card_w - Inches(0.9),
        Inches(0.5),
        "+44 (0) 1344 231 119",
        font=BODY_FONT,
        size=16,
        color=SLATE_SOFT,
    )
    add_text(
        slide,
        card_x + Inches(0.45),
        card_y + Inches(1.85),
        card_w - Inches(0.9),
        Inches(0.6),
        "12 Longshot Lane, Bracknell, Berkshire, RG12 1RL",
        font=BODY_FONT,
        size=12,
        color=SLATE,
        line_spacing=1.3,
    )

    # Legal card (right)
    legal_x = Inches(6.8)
    legal_y = Inches(4.3)
    legal_w = Inches(5.93)
    legal_h = Inches(2.6)
    add_rounded_rect(slide, legal_x, legal_y, legal_w, legal_h, INK_RAISED, WHITE_10, radius=0.05)
    add_text(
        slide,
        legal_x + Inches(0.45),
        legal_y + Inches(0.35),
        Inches(3),
        Inches(0.35),
        "THE LEGAL STUFF",
        font=MONO_FONT,
        size=10,
        color=VIOLET,
        tracking=260,
    )
    add_text(
        slide,
        legal_x + Inches(0.45),
        legal_y + Inches(0.8),
        legal_w - Inches(0.9),
        Inches(1.6),
        "50% deposit up front, 50% due before launch. Maintenance is "
        "£120/month paid in advance by Direct Debit, starting the month "
        "after go-live. All prices exclude VAT. T&Cs on signing.",
        font=BODY_FONT,
        size=12,
        color=SLATE_SOFT,
        line_spacing=1.5,
    )
    add_text(
        slide,
        legal_x + Inches(0.45),
        legal_y + Inches(2.05),
        legal_w - Inches(0.9),
        Inches(0.4),
        "WEBGRO LTD  ·  REG 10889889  ·  PART OF BROADBRIDGE GROUP",
        font=MONO_FONT,
        size=9,
        color=SLATE,
        tracking=220,
    )

    add_footer(slide)
    add_page_number(slide, n, total)


# -----------------------------------------------------------------------------
# Main
# -----------------------------------------------------------------------------


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Count slides up front so page numbers read "N / 13"
    # 1 cover + 4 section dividers + 8 content = 13
    TOTAL = 13

    # 1. Cover
    build_cover(prs)

    # 2. Section 1 divider — Get to know us
    build_section_divider(
        prs,
        "01.",
        "Get to know us.",
        "The studio, the record, and the awards before we talk about you.",
        BLUE,
        2,
        TOTAL,
    )

    # 3. Hello / who we are
    build_hello(prs, 3, TOTAL)

    # 4. Six capabilities
    build_capabilities(prs, 4, TOTAL)

    # 5. Awards
    build_awards(prs, 5, TOTAL)

    # 6. Section 2 divider — Work
    build_section_divider(
        prs,
        "02.",
        "Selected work.",
        "Three WordPress service businesses, rebuilt, launched, and kept current.",
        VIOLET,
        6,
        TOTAL,
    )

    # 7-9. Case studies — WordPress service-business examples, chosen to
    # show the same shape of work as Ithaca Roofing's brief rather than
    # leading with eCommerce.
    build_case_study(
        prs,
        7,
        TOTAL,
        client="Little Muddy Boots",
        url="littlemuddyboots.co.uk",
        tagline="A mobile-first WordPress redesign for a service business, with a custom postcode-search tool so visitors know instantly whether the service covers their area.",
        services="WordPress redesign, custom postcode search, mobile-first UX",
        image_path=IMG_DIR / "little-muddy-boots.jpg",
        accent=BLUE,
        year="2025",
    )
    build_case_study(
        prs,
        8,
        TOTAL,
        client="Space 4 U Self Storage",
        url="space4uselfstorage.co.uk",
        tagline="A WordPress redesign for a regional self-storage operator. Clarity, speed, and a direct path from visitor to unit enquiry.",
        services="WordPress redesign, conversion-led UX, ongoing maintenance",
        image_path=IMG_DIR / "space4u.jpg",
        accent=VIOLET,
        year="2025",
    )
    build_case_study(
        prs,
        9,
        TOTAL,
        client="Paragon Freight",
        url="paragonfreight.co.uk",
        tagline="A B2B WordPress rebuild for a freight operator, with a live currency selector and enquiry flow shaped around international buyers.",
        services="WordPress build, currency selector, enquiry flow, retainer",
        image_path=IMG_DIR / "paragon-freight.jpg",
        accent=TEAL,
        year="2025 → present",
    )

    # 10. Section 3 divider — Proposal
    build_section_divider(
        prs,
        "03.",
        "What we'll do for you.",
        "From Ithaca Roofing's brief, translated into scope and price.",
        TEAL,
        10,
        TOTAL,
    )

    # 11. What we provide
    build_what_we_provide(prs, 11, TOTAL)

    # 12. Investment
    build_investment(prs, 12, TOTAL)

    # 13. Ready to gro / contact
    build_ready(prs, 13, TOTAL)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
